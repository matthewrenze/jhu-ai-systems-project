# Import libraries
import os
import pptx
from PIL import Image
from io import BytesIO

# Define the input and output file paths
root_folder_path = "C:\\Users\\Matthew\\Dropbox\\Personal\\School\\JHU\\Systems\\Project"
input_folder_path = root_folder_path + "\\Data\\PowerPoint\\Intro to Data for Data Science"
output_folder_path = root_folder_path + "\\Data\\Markdown\\Intro to Data for Data Science"

def get_sections(presentation):
    extLst = '{http://schemas.openxmlformats.org/presentationml/2006/main}extLst'
    ext = '{http://schemas.openxmlformats.org/presentationml/2006/main}ext'
    uri = '{521415D9-36F7-43E2-AB2F-B90AF26B5E84}'
    element = presentation.element.find(f"./{extLst}/{ext}[@uri='{uri}']")

    if element is None:
        slides = presentation.slides
        slide_ids = [slide.slide_id for slide in slides]
        return [None], {None: slide_ids}

    section_list = element[0]

    sections = []
    section_map = {}
    for section in section_list:
        section_name = section.attrib['name']
        sections.append(section_name)
        slide_ids = []
        slide_list = section[0]
        for slide_element in slide_list:
            slide_id = int(slide_element.attrib['id'])
            slide_ids.append(slide_id)
        section_map[section_name] = slide_ids

    return sections, section_map

def get_slide_map(presentation):
    slide_map = {}
    for slide in presentation.slides:
        slide_id = slide.slide_id
        slide_map[slide_id] = slide
    return slide_map

def get_presentation_title(section_index, section_name):
    if section_name is not None:
        return f"{section_index} - {section_name}"
    else:
        return presentation_name

def is_hidden(slide):
    return slide._element.get("show") == "0"

def get_format(slide):
    power_point_layout_map = {
        "Title Slide": "title-slide",
        "Section Header": "section-slide",
        "Section Slide": "section-slide",
        "Title Only": "title-only",
        "Title and Content": "title-one-content-left",
        "Two Content": "title-two-content-left",
        "Comparison": "title-two-content-comparison"
    }

    my_layout_map = {
        "Title One Content": "title-one-content",
        "Title One Content Left": "title-one-content-left",
        "Title Two Content": "title-two-content",
        "Title Two Content Left": "title-two-content-left",
        "Title Two Content Inner": "two-content-inner",
        "Title Two Content Left Center": "title-two-content-left-center",
        "Title Two Content Center Left": "title-two-content-center-left",
        "Title Two Content Comparison": "title-two-content-comparison",
        "Title Three Content": "title-three-content",
        "Title Four Content (Vertical)": "title-four-content",
        "Title Four Content Two Row": "title-four-content-two-row",
        "Title Five Content": "title-five-content",
        "Title Six Content": "title-six-content",
        "Title Seven Content": "title-seven-content",
        "Title Eight Content": "title-eight-content",
        "One Pane": "one-pane",
        "Two Pane": "two-pane",
        "Two Pane Golden": "two-pane-golden",
        "Three Pane": "three-pane",
        "Four Pane": "four-pane"
    }

    pluralsight_map = {
        "Two Text": "title-two-content-left",
        "Three Text": "title-three-content",
        "Four Text": "title-four-content",
        "Five Text": "title-five-content",
        "Six Text": "title-six-content",
        "Two Item": "title-two-content",
        "Three Item": "title-three-content",
        "Four Item": "title-four-content",
        "Five Item": "title-five-content",
        "Six Item": "title-six-content",
    }

    layout_map = {**power_point_layout_map, **my_layout_map, **pluralsight_map}

    layout_name = slide.slide_layout.name

    if layout_name in layout_map:
        return layout_map[layout_name]
    else:
        return None

def is_title(shape):
    if shape.is_placeholder:
        if shape.placeholder_format.type._member_name in ["TITLE", "CENTER_TITLE", "SUBTITLE"]:
            return True
    return False

def format_title(text):
    return f"# {text}"

def get_title(slide):
    title = None
    if slide.shapes.title:
        title = format_title(slide.shapes.title.text)
    else:
        for shape in slide.shapes:
            if is_title(shape):
                title = format_title(shape.text_frame.text)
    title = clean_text(title)
    return title

def has_text(shape):
    return shape.has_text_frame

def clean_text(text):
    if text is None:
        return None

    text = text.replace("\x0b", "\n") \
        .replace("’", "'") \
        .replace("“", "\"") \
        .replace("”", "\"") \
        .replace("…", "...") \
        .replace("–", "-") \
        .replace("✔", "x")
    text = text.strip()
    return text

def get_text(shape):
    text = shape.text_frame.text
    text = clean_text(text)
    return text

def is_image(shape):
    is_image = True if shape.shape_type._member_name == "PICTURE" else False
    may_have_image = shape.shape_type._member_name == "OBJECT" \
        or (shape.shape_type._member_name == "PLACEHOLDER"
            and shape.has_text_frame == False
            and (shape.placeholder_format.type._member_name == "PICTURE"
                or shape.placeholder_format.type._member_name == "OBJECT"))
    if may_have_image:
        is_image = True if shape.image else False
    return is_image

def is_background_image(shape, shape_index):
    if shape.shape_type._member_name == "PICTURE" \
            and shape_index == 0:
        return True

def get_image_file_name(shape):
    slide_id = slide.slide_id
    shape_id = shape.shape_id
    extension = shape.image.filename.split(".")[-1]
    image_file_name = f"{slide_id}-{shape_id}.{extension}"
    return image_file_name

def crop_image(image):
    crop_left = shape.crop_left
    crop_right = shape.crop_right
    crop_top = shape.crop_top
    crop_bottom = shape.crop_bottom
    image_bytes = BytesIO(image)
    image = Image.open(image_bytes)
    image_width, image_height = image.size
    new_left = int(image_width * crop_left)
    new_right = int(image_width - image_width * crop_right)
    new_top = int(image_height * crop_top)
    new_bottom = int(image_height - image_height * crop_bottom)
    image = image.crop((new_left, new_top, new_right, new_bottom))
    return image

def resize_image(image):
    max_width = 4000
    max_height = 2250
    image_width, image_height = image.size
    if image_width > max_width or image_height > max_height:
        new_width = max_width
        new_height = int(new_width * image_height / image_width)
        if new_height > max_height:
            new_height = max_height
            new_width = int(new_height * image_width / image_height)
        image = image.resize((new_width, new_height))
    return image

def save_image(image, image_folder_path, file_name):
    file_name = os.path.splitext(file_name)[0] + ".png"
    file_path = image_folder_path + "/" + file_name
    image.save(file_path)

def get_image_alt_text(shape):
    pic_element = shape._element.xpath('.//p:nvPicPr')[0]
    cNvPr_element = pic_element.xpath('.//p:cNvPr')[0]
    alt_text = cNvPr_element.get('descr')
    alt_text = alt_text if alt_text is not None else ""
    alt_text = alt_text.replace("\n", "")
    alt_text = alt_text.replace("Description automatically generated", "")
    alt_text = alt_text.strip()
    alt_text = " " + alt_text if alt_text != "" else ""
    return alt_text

def has_table(shape):
    return shape.shape_type._member_name == "TABLE"

def get_table(shape):
    table = shape.table
    table_text = ""
    for r in range(0, len(table.rows)):
        row_text = "|"
        for c in range(0, len(table.columns)):
            cell = table.cell(r, c)
            cell_text = cell.text
            row_text += f" {cell_text} |"
        table_text += row_text + "\n"
        if r == 0:
            divider_text = "| "
            for c in range(0, len(table.columns)):
                divider_text += " --- |"
            table_text += divider_text + "\n"

    return table_text

def is_footer(shapes, shape):
    if len(shapes) < 2:
        return False

    if shape != shapes[-1]:
        return False

    if not shape.has_text_frame:
        return False

    if shape.shape_type._member_name == "AUTO_SHAPE":
        return False

    if shape.text_frame is None \
        or shape.text_frame.paragraphs is None \
        or shape.text_frame.paragraphs[0] is None \
        or len(shape.text_frame.paragraphs[0].runs) == 0 \
        or shape.text_frame.paragraphs[0].runs[0] is None \
        or shape.text_frame.paragraphs[0].runs[0].font is None \
        or shape.text_frame.paragraphs[0].runs[0].font.size is None \
        or shape.text_frame.paragraphs[0].runs[0].font.size.pt is None:
        return False

    font_size = shape.text_frame.paragraphs[0].runs[0].font.size.pt
    if font_size > 12:
        return False

    text = get_text(shape)

    if text is not None and text.find("\n") == -1:
        return True

    return False

def has_notes(slide):
    return slide.has_notes_slide

def get_notes(slide):
    notes = slide.notes_slide.notes_text_frame.text
    notes = clean_text(notes)
    return notes

# Print an inital status update
print("Beginning all conversions...")

# Get all files with ".pptx" extensions in all subfolders
input_file_paths = []
for root, dirs, file_names in os.walk(input_folder_path):
    for file_name in file_names:
        if file_name.endswith('.pptx'):
            input_file_paths.append(os.path.join(root, file_name))

# Convert each presentation to markdown
for input_file_path in input_file_paths:

    # Print a status update
    print(f"Converting presentation {os.path.basename(input_file_path)}...")

    # Create the output file paths
    input_file_name = os.path.basename(input_file_path)
    presentation_name = input_file_name.replace(".pptx", "")
    output_subfolder_path = output_folder_path + "\\" + presentation_name
    image_folder_path = "images"

    # Create output folder if it does not exist
    if not os.path.exists(output_folder_path):
        os.makedirs(output_folder_path)

    # Create the image folder if it doesn't exist
    if not os.path.exists(image_folder_path):
        os.makedirs(image_folder_path)

    # Load the PowerPoint file
    presentation = pptx.Presentation(input_file_path)

    # Get the sections and slide map
    sections, section_map = get_sections(presentation)
    slide_map = get_slide_map(presentation)
    section_index = 1

    # Convert each section
    for section in sections:

        # Print a status update
        print(f"Converting section {section_index} - {section}...")

        # Get the section details
        section_name = section
        slide_ids = section_map[section_name]
        presentation_title = get_presentation_title(section_index, section_name)

        # Create markdown header
        text = "---\n"
        text += "marp: true\n"
        text += f"title: {presentation_title}\n"
        text += "theme: template\n"
        text += "---\n\n"

        for slide_id in slide_ids:

            slide = slide_map[slide_id]

            if is_hidden(slide):
                continue

            format = get_format(slide)
            if format:
                text += f"<!-- _class: {format} -->\n\n"

            title = get_title(slide)
            if title is not None and title != "":
                text += f"{title}\n\n"

            for shape in slide.shapes:
                shape_index = slide.shapes.index(shape)

                if is_title(shape) and title is not None:
                    continue

                if has_text(shape):
                    if is_footer(slide.shapes, shape):
                        text += f"<!-- _footer: {get_text(shape)} -->\n"
                    else:
                        text += f"{get_text(shape)}\n"

                if is_image(shape):
                    image_file_name = get_image_file_name(shape)
                    alt_text = get_image_alt_text(shape)
                    cropped_image = crop_image(shape.image.blob)
                    resized_image = resize_image(cropped_image)
                    save_image(resized_image, image_folder_path, image_file_name)
                    if is_background_image(shape, shape_index):
                        text += f"![bg contain{alt_text}]({image_folder_path}/{image_file_name})\n"
                    else:
                        text += f"![image{alt_text}]({image_folder_path}/{image_file_name})\n"

                if has_table(shape):
                    text += f"{get_table(shape)}\n"

                text += "\n"

            if has_notes(slide):
                text += f"<!--\n{get_notes(slide)}\n-->\n"

            if slide.slide_id != slide_ids[-1]:
                text += "\n---\n\n"

            # Print a status update
            print(f"Converted slide {slide.slide_id} - {title}")

        # Clean text of any un-encodable characters
        text = clean_text(text)

        # Create the output file path
        output_file_name = presentation_title + ".md"
        output_file_path = output_subfolder_path + "\\" + output_file_name

        # Save the markdown to a file
        with open(output_file_path, "w") as text_file:
            text_file.write(text)

        # Print a status update
        print(f"Saved section {output_file_name}\n")

        section_index += 1

# Print a final status update
print("All conversions complete")

